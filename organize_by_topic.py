import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from collections import Counter
import numpy as np
import argparse
import logging
import os
from pathlib import Path
from itertools import product
from openai import OpenAI
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
import io

# Set up logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logging.getLogger("httpx").setLevel(logging.WARNING)
logger = logging.getLogger(__name__)

# Set global matplotlib font size
plt.rcParams.update({'font.size': 18})

# Load environment variables
load_dotenv()

# Set up the plotting style (same as plot_analysis.py)
plt.style.use('seaborn-v0_8')
sns.set_palette("husl")

# ALCF primer for context
ALCF_PRIMER = """The Argonne Leadership Computing Facility (ALCF) is a U.S. Department of Energy (DOE) Office of Science user facility that provides supercomputing resources and expertise to the scientific and engineering community. ALCF's mission is to accelerate major scientific discoveries and engineering breakthroughs for humanity by designing and providing world-leading computing facilities in partnership with the computational science community. The facility supports a wide range of research areas including climate science, materials science, physics, chemistry, and more. ALCF provides access to high-performance computing systems, scientific software, and user support services to enable researchers to advance their work."""

def parse_args():
    parser = argparse.ArgumentParser(description='Organize analysis results by topic')
    parser.add_argument('-i', '--input', type=str, default='analysis_results.csv',
                      help='Input CSV file path (default: analysis_results.csv)')
    parser.add_argument('-o', '--output-dir', type=str, default='topic_analysis',
                      help='Output directory for topic analysis (default: topic_analysis)')
    return parser.parse_args()

def plot_comma_separated_histogram(series, title, figsize=(10, 6)):
    """Create a histogram for comma-separated values."""
    all_values = []
    for value in series:
        if pd.notna(value):
            all_values.extend([v.strip() for v in str(value).split(',')])
    
    value_counts = Counter(all_values)
    total_count = sum(value_counts.values())
    
    # Convert to percentages
    percentages = {k: (v / total_count) * 100 for k, v in value_counts.items()}
    
    plt.figure(figsize=figsize)
    bars = plt.bar(percentages.keys(), percentages.values())
    
    # Add total count to title
    title_with_count = f"{title} (Total: {total_count})"
    plt.title(title_with_count, pad=20, fontsize=20)
    plt.xticks(rotation=45, ha='right', fontsize=20)
    plt.yticks(fontsize=20)
    plt.xlabel('Categories', fontsize=20)
    plt.ylabel('Percentage (%)', fontsize=20)
    plt.tight_layout()
    
    for bar in bars:
        height = bar.get_height()
        plt.text(bar.get_x() + bar.get_width()/2., height,
                f'{height:.1f}%',
                ha='center', va='bottom')
    
    return plt.gcf()

def plot_regular_histogram(series, title, figsize=(10, 6)):
    """Create a regular histogram."""
    plt.figure(figsize=figsize)
    value_counts = series.value_counts()
    total_count = len(series.dropna())
    
    # Convert to percentages
    percentages = (value_counts / total_count) * 100
    
    bars = plt.bar(percentages.index, percentages.values)
    
    # Add total count to title
    title_with_count = f"{title} (Total: {total_count})"
    plt.title(title_with_count, pad=20, fontsize=20)
    plt.xticks(rotation=45, ha='right', fontsize=20)
    plt.yticks(fontsize=20)
    plt.xlabel('Categories', fontsize=20)
    plt.ylabel('Percentage (%)', fontsize=20)
    plt.tight_layout()
    
    for bar in bars:
        height = bar.get_height()
        plt.text(bar.get_x() + bar.get_width()/2., height,
                f'{height:.1f}%',
                ha='center', va='bottom')
    
    return plt.gcf()

def plot_2d_histogram(df, col1, col2, title, figsize=(12, 8)):
    """Create a 2D histogram between two columns."""
    col1_values = set()
    col2_values = set()
    
    for val1, val2 in zip(df[col1], df[col2]):
        if pd.notna(val1) and pd.notna(val2):
            col1_values.update([v.strip() for v in str(val1).split(',')])
            col2_values.update([v.strip() for v in str(val2).split(',')])
    
    col1_values = sorted(list(col1_values))
    col2_values = sorted(list(col2_values))
    
    hist_data = np.zeros((len(col2_values), len(col1_values)))
    
    for val1, val2 in zip(df[col1], df[col2]):
        if pd.notna(val1) and pd.notna(val2):
            val1_list = [v.strip() for v in str(val1).split(',')]
            val2_list = [v.strip() for v in str(val2).split(',')]
            
            for v1, v2 in product(val1_list, val2_list):
                i = col1_values.index(v1)
                j = col2_values.index(v2)
                hist_data[j, i] += 1
    
    # Convert to percentages
    total_count = np.sum(hist_data)
    hist_data_percent = (hist_data / total_count) * 100
    
    plt.figure(figsize=figsize)
    sns.heatmap(hist_data_percent, 
                xticklabels=col1_values,
                yticklabels=col2_values,
                cmap='YlOrRd',
                annot=True,
                fmt='.1f',
                cbar_kws={'label': 'Percentage (%)'})
    
    # Add total count to title
    title_with_count = f"{title} (Total: {total_count})"
    plt.title(title_with_count, pad=20, fontsize=20)
    plt.xlabel(col1.title(), fontsize=20)
    plt.ylabel(col2.title(), fontsize=20)
    plt.xticks(rotation=45, ha='right', fontsize=20)
    plt.yticks(rotation=0, fontsize=20)
    plt.tight_layout()
    
    return plt.gcf()

def analyze_feedback(client, feedback_data):
    """Analyze the feedback using the LLM and generate a summary and action items."""
    feedback_text = "\n\n".join([
        f"Response {i+1}: {row['user_response']}"
        for i, row in feedback_data.iterrows()
    ])

    prompt = f"""{ALCF_PRIMER}

Below are user feedback responses about ALCF's software development environment:

{feedback_text}

Please provide:
1. A concise summary of the main themes and concerns expressed in the feedback calling out specific examples from the feedback. These should be specific things like "the documentation on topic X was not helpful" or "the filesystem on Polaris is not reliable". This should not be an itemized list, but a summary paragraph.
2. The top three actionable items that ALCF should consider to address these concerns. Again, these should be specific items and not generic things like "improve documentation" or "improve support".

Format your response as follows:

SUMMARY:
[Your summary here]

TOP THREE ACTION ITEMS:
1. [First action item]
2. [Second action item]
3. [Third action item]
"""

    try:
        response = client.chat.completions.create(
            model=os.environ.get("OPENAI_MODEL_NAME", "argo:gpt-4"),
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        logger.error(f"Error analyzing feedback: {str(e)}")
        return "ERROR: Failed to analyze feedback"

def create_powerpoint(topic, markdown_content, plots_dir):
    """Create a PowerPoint presentation with the markdown content and plots."""
    prs = Presentation()
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = f"Feedback Analysis: {topic.title()}"
    subtitle.text = "ALCF User Survey Results"
    
    # Set font size for title and subtitle
    title.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.size = Pt(12)
    
    # Summary slide
    summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = summary_slide.shapes.title
    content = summary_slide.placeholders[1]
    title.text = "Summary"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    
    # Extract summary from markdown content
    summary = markdown_content.split("TOP THREE ACTION ITEMS:")[0].replace("SUMMARY:", "").strip()
    content.text = summary
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(12)
    
    # Action items slide
    action_slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = action_slide.shapes.title
    content = action_slide.placeholders[1]
    title.text = "Top Three Action Items"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    
    # Extract action items from markdown content
    try:
        action_items = markdown_content.split("TOP THREE ACTION ITEMS:")[1].strip()
    except:
        logger.error(f"Error extracting action items from markdown content: {markdown_content}")
        raise
    content.text = action_items
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(12)
    
    # Add plot slides
    for plot_file in plots_dir.glob("*.png"):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = plot_file.stem.replace("_", " ").title()
        title.text_frame.paragraphs[0].font.size = Pt(24)
        
        # Add the plot
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        slide.shapes.add_picture(str(plot_file), left, top, width, height)
    
    return prs

def main():
    args = parse_args()
    
    # Create output directory
    output_dir = Path(args.output_dir)
    output_dir.mkdir(exist_ok=True)
    logger.info(f"Created output directory: {output_dir}")
    
    # Read the CSV file
    logger.info(f"Reading input file: {args.input}")
    df = pd.read_csv(args.input)
    
    # Normalize classification columns to lowercase to avoid case inconsistencies in plots
    classification_columns = ['sentiment', 'topic', 'machine', 'feedback', 'emotion', 
                            'actionability', 'specificity', 'ticket_status', 'software_topic']
    for col in classification_columns:
        if col in df.columns:
            df[col] = df[col].astype(str).str.lower()
    logger.info("Normalized classification columns to lowercase")
    
    # Initialize OpenAI client
    client = OpenAI(
        api_key=os.environ.get("OPENAI_API_KEY", "local-key"),
        base_url=os.environ.get("OPENAI_API_BASE", "http://localhost:44497/v1")
    )
    
    # Process ALL data first (aggregate analysis)
    logger.info("Processing aggregate analysis for all data")
    overall_dir = output_dir / "overall"
    overall_dir.mkdir(exist_ok=True)
    
    # Save all responses
    df.to_csv(overall_dir / 'all_responses.csv')
    
    # Filter negative responses for overall analysis
    overall_negative_responses = df[df['sentiment'] == 'negative']
    overall_negative_responses.to_csv(overall_dir / 'negative_responses.csv')
    
    # Create plots directory for overall analysis
    overall_plots_dir = overall_dir / 'plots'
    overall_plots_dir.mkdir(exist_ok=True)
    
    # Generate plots from ALL responses (aggregate)
    columns_to_plot = {
        'sentiment': 'regular',
        'machine': 'comma_separated',
        'feedback': 'regular',
        'emotion': 'regular',
        'actionability': 'regular',
        'specificity': 'regular',
        'ticket_status': 'regular',
        'software_topic': 'comma_separated'
    }
    
    for col, plot_type in columns_to_plot.items():
        title = f'{col.title()} Distribution - Overall'
        if plot_type == 'comma_separated':
            fig = plot_comma_separated_histogram(df[col], title)
        else:
            fig = plot_regular_histogram(df[col], title)
        
        fig.savefig(overall_plots_dir / f'{col}_distribution.png', bbox_inches='tight', dpi=300)
        plt.close(fig)
    
    # Generate 2D histograms from ALL responses (aggregate)
    correlations = [
        ('machine', 'emotion', 'Machine vs Emotion'),
        ('feedback', 'emotion', 'Feedback Type vs Emotion'),
        ('software_topic', 'emotion', 'Software Topic vs Emotion'),
        ('specificity', 'actionability', 'Specificity vs Actionability'),
        ('ticket_status', 'emotion', 'Ticket Status vs Emotion'),
        ('sentiment', 'machine', 'Sentiment vs Machine'),
        ('sentiment', 'feedback', 'Sentiment vs Feedback Type'),
        ('sentiment', 'emotion', 'Sentiment vs Emotion'),
        ('sentiment', 'actionability', 'Sentiment vs Actionability'),
        ('sentiment', 'specificity', 'Sentiment vs Specificity'),
        ('sentiment', 'ticket_status', 'Sentiment vs Ticket Status'),
        ('sentiment', 'software_topic', 'Sentiment vs Software Topic')
    ]
    
    for col1, col2, title_prefix in correlations:
        title = f'{title_prefix} Distribution - Overall'
        fig = plot_2d_histogram(df, col1, col2, title)
        fig.savefig(overall_plots_dir / f'{col1}_{col2}_2d_histogram.png', bbox_inches='tight', dpi=300)
        plt.close(fig)
    
    if not overall_negative_responses.empty:
        # Analyze negative feedback for overall data
        overall_analysis_result = analyze_feedback(client, overall_negative_responses)
        # catch connection errors
        if "500 Internal Server Error" in overall_analysis_result:
            raise Exception(f"Error analyzing feedback: {overall_analysis_result}")
        
        # Save markdown file for overall analysis
        with open(overall_dir / 'negative_feedback_analysis.md', 'w') as f:
            f.write("# Overall Feedback Analysis Results\n\n")
            f.write(overall_analysis_result)
        
        # Create PowerPoint presentation for overall analysis
        prs = create_powerpoint("Overall", overall_analysis_result, overall_plots_dir)
        prs.save(overall_dir / 'overall_analysis.pptx')
        
        logger.info("Completed overall aggregate analysis")
    else:
        logger.info("No negative responses found in overall data")

    # Get all unique topics
    all_topics = set()
    for topics in df['topic'].dropna():
        all_topics.update(t.strip() for t in str(topics).split(','))
    
    # Process each topic
    for topic in all_topics:
        logger.info(f"Processing topic: {topic}")
        topic_dir = output_dir / topic.replace(" ", "_")
        topic_dir.mkdir(exist_ok=True)
        
        # Filter responses for this topic
        topic_responses = df[df['topic'].str.contains(topic, na=False)]
        topic_responses.to_csv(topic_dir / 'all_responses.csv')
        
        # Filter negative responses
        negative_responses = topic_responses[topic_responses['sentiment'] == 'negative']
        negative_responses.to_csv(topic_dir / 'negative_responses.csv')
        
        # Create plots directory (now always, not just for negative)
        plots_dir = topic_dir / 'plots'
        plots_dir.mkdir(exist_ok=True)
        
        # Generate plots from ALL responses for this topic
        columns_to_plot = {
            'sentiment': 'regular',
            'machine': 'comma_separated',
            'feedback': 'regular',
            'emotion': 'regular',
            'actionability': 'regular',
            'specificity': 'regular',
            'ticket_status': 'regular',
            'software_topic': 'comma_separated'
        }
        
        for col, plot_type in columns_to_plot.items():
            title = f'{col.title()} Distribution - {topic.title()}'
            if plot_type == 'comma_separated':
                fig = plot_comma_separated_histogram(topic_responses[col], title)
            else:
                fig = plot_regular_histogram(topic_responses[col], title)
            
            fig.savefig(plots_dir / f'{col}_distribution.png', bbox_inches='tight', dpi=300)
            plt.close(fig)
        
        # Generate 2D histograms from ALL responses for this topic
        correlations = [
            ('machine', 'emotion', 'Machine vs Emotion'),
            ('feedback', 'emotion', 'Feedback Type vs Emotion'),
            ('software_topic', 'emotion', 'Software Topic vs Emotion'),
            ('specificity', 'actionability', 'Specificity vs Actionability'),
            ('ticket_status', 'emotion', 'Ticket Status vs Emotion'),
            ('sentiment', 'machine', 'Sentiment vs Machine'),
            ('sentiment', 'feedback', 'Sentiment vs Feedback Type'),
            ('sentiment', 'emotion', 'Sentiment vs Emotion'),
            ('sentiment', 'actionability', 'Sentiment vs Actionability'),
            ('sentiment', 'specificity', 'Sentiment vs Specificity'),
            ('sentiment', 'ticket_status', 'Sentiment vs Ticket Status'),
            ('sentiment', 'software_topic', 'Sentiment vs Software Topic')
        ]
        
        for col1, col2, title_prefix in correlations:
            title = f'{title_prefix} Distribution - {topic.title()}'
            fig = plot_2d_histogram(topic_responses, col1, col2, title)
            fig.savefig(plots_dir / f'{col1}_{col2}_2d_histogram.png', bbox_inches='tight', dpi=300)
            plt.close(fig)
        
        if not negative_responses.empty:
            # Analyze negative feedback
            analysis_result = analyze_feedback(client, negative_responses)
            # catch connection errors
            if "500 Internal Server Error" in analysis_result:
                # logger.error(f"Error analyzing feedback: {analysis_result}")
                raise Exception(f"Error analyzing feedback: {analysis_result}")
            
            # Save markdown file
            with open(topic_dir / 'negative_feedback_analysis.md', 'w') as f:
                f.write("# Feedback Analysis Results\n\n")
                f.write(analysis_result)
            
            # Create PowerPoint presentation
            prs = create_powerpoint(topic, analysis_result, plots_dir)
            prs.save(topic_dir / f'{topic.replace(" ", "_")}_analysis.pptx')
            
            logger.info(f"Completed processing for topic: {topic}")
        else:
            logger.info(f"No negative responses found for topic: {topic}")

if __name__ == '__main__':
    main() 